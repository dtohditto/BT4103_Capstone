import os
import pandas as pd
import numpy as np
import streamlit as st
import plotly.express as px
import io, zipfile, math, textwrap
import plotly.io as pio
from datetime import datetime

# Exports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader

import CSVCuration

st.set_page_config(page_title="EE Analytics Dashboard", layout="wide")

# ------------------------------
# Data loading & preparation
# ------------------------------
@st.cache_data
def load_data(path: str = "dashboard_curated_v2.csv") -> pd.DataFrame:
    """Load curated CSV and perform light, safe parsing for the app."""
    df = pd.read_csv(path)

    # Parse key dates (coerce errors to NaT)
    for c in ["Programme Start Date", "Programme End Date", "Run_Month"]:
        if c in df.columns:
            df[c] = pd.to_datetime(df[c], errors="coerce")

    # Friendly label for charts (if not already present)
    if "Run_Month" in df.columns:
        df["Run_Month_Label"] = df["Run_Month"].dt.strftime("%Y-%m")

    # Derive Age_Group once 
    if "Age" in df.columns:
        age_num = pd.to_numeric(df["Age"], errors="coerce")
        bins    = [0, 34, 44, 54, 64, 200]
        labels  = ["<35", "35–44", "45–54", "55–64", "65+"]
        df["Age_Group"] = pd.cut(age_num, bins=bins, labels=labels, right=True).astype("string")

    # Light category cleanup for common slicers
    cat_cols = [
        "Application Status", "Applicant Type", "Primary Category",
        "Secondary Category", "Seniority", "Gender", "Country Of Residence",
        "Truncated Programme Name", "Domain"
    ]
    for c in cat_cols:
        if c in df.columns:
            df[c] = df[c].astype("string").str.strip()

    return df


# Allow user to upload a newer CSV (optional)
new_uploaded_programme = st.sidebar.file_uploader("Upload a new programme CSV (optional)", type=["xlsm", "xlsx", "csv"])
new_uploaded_cost = st.sidebar.file_uploader("Upload a new cost CSV (optional)", type=["xlsm", "xlsx", "csv"])
# call the function and request CSV bytes
res = CSVCuration.curate_programme_and_cost_data(new_uploaded_programme, new_uploaded_cost, return_csv_bytes=True)

# res may be None, bytes, a DataFrame, or (DataFrame, bytes)
csv_bytes = None
df_curated = None
if res is None:
    csv_bytes = None
elif isinstance(res, tuple):
    df_curated, csv_bytes = res
elif isinstance(res, (bytes, bytearray)):
    csv_bytes = bytes(res)
elif isinstance(res, pd.DataFrame):
    df_curated = res
    # produce bytes so the download button can work
    try:
        csv_bytes = df_curated.to_csv(index=False).encode("utf-8-sig")
    except Exception:
        csv_bytes = None

if csv_bytes is not None:
    st.sidebar.download_button(
        "Download curated CSV",
        data=csv_bytes,
        file_name="dashboard_curated.csv",
        mime="text/csv"
    )
else:
    # No curated CSV available yet; show a helpful note instead of an invalid download button
    st.sidebar.info("No curated CSV available — upload both programme and cost files to generate a curated CSV.")
uploaded = st.sidebar.file_uploader("Upload a curated CSV (optional)", type=["csv"])
data_path = uploaded if uploaded is not None else "dashboard_curated_v2.csv"

# Load the full dataset once; df is a working copy for filters
df_full = load_data(data_path)
df = df_full.copy()

if df.empty:
    st.info("No data loaded. Please upload a CSV or place dashboard_curated.csv next to this script.")
    st.stop()

st.title("Executive Education Analytics Dashboard")

# ------------------------------
# Global constants & helpers
# ------------------------------
UNKNOWN_LIKE = {
    "unknown", "unspecified", "not specified", "not provided", "not available",
    "n/a", "na", "null", "none", "-", "", 
}

def _safe_key(label: str, suffix: str) -> str:
    return f"{label}_{suffix}".replace(" ", "_").lower()

# UI label ↔ column mapping
COL_MAP = {
    "Pri Category": "Primary Category",
    "Sec Category": "Secondary Category",
    "Country": "Country Of Residence",
    "Application Status": "Application Status",
    "Applicant Type": "Applicant Type",
    "Seniority": "Seniority",
    "Domain": "Domain",
}
UI_FILTER_LABELS = ["Application Status", "Applicant Type", "Pri Category",
                    "Sec Category", "Country", "Seniority", "Domain"]

def _col_from_label(label: str) -> str:
    return COL_MAP.get(label, label)

def multiselect_with_all_button(label: str, df_source: pd.DataFrame, default_all: bool = True):
    col = _col_from_label(label)
    raw = df_source.get(col, pd.Series([], dtype="object")).copy()
    # If the column is categorical, convert to string first so fillna can add 'Unknown'
    if pd.api.types.is_categorical_dtype(raw):
        s = raw.astype("string").fillna("Unknown")
    else:
        s = raw.fillna("Unknown")
    options = sorted(s.unique().tolist())

    ms_key = _safe_key(label, "multi")
    if ms_key not in st.session_state:
        st.session_state[ms_key] = options[:] if default_all else []

    current = [v for v in st.session_state[ms_key] if v in options]
    if current != st.session_state[ms_key]:
        st.session_state[ms_key] = current

    if st.button(f"Select all {label}", key=_safe_key(label, "btn_all")):
        st.session_state[ms_key] = options[:]

    st.multiselect(label, options, key=ms_key)
    return st.session_state[ms_key]

def apply_filter(series: pd.Series, selected: list[str]) -> pd.Series:
    if selected is None:
        return pd.Series(True, index=series.index)
    # Avoid fillna on Categorical dtype (raises if the fill value is not a category)
    if pd.api.types.is_categorical_dtype(series):
        s = series.astype("string").fillna("Unknown")
    else:
        s = series.fillna("Unknown")
    all_opts = set(s.unique().tolist())
    sel_set  = set(selected or [])
    if not selected or sel_set == all_opts:
        return pd.Series(True, index=series.index)
    return s.isin(selected)

# ---- Unknown/Missing helpers ----

def _norm_str(s: pd.Series) -> pd.Series:
    return s.astype("string").str.strip().str.lower()

def coalesce_unknown(series: pd.Series) -> pd.Series:
    s = series.astype("string").str.strip()
    s_norm = s.str.lower()
    mask_u = s.isna() | (s == "") | s_norm.isin(UNKNOWN_LIKE)
    return s.mask(mask_u, "Unknown")

def filter_unknown_no_ui(df_in: pd.DataFrame, column: str, include_unknown: bool) -> pd.DataFrame:
    if column not in df_in.columns:
        return df_in
    s_norm = df_in[column].astype("string").str.strip().str.lower()
    mask_unknown_or_missing = df_in[column].isna() | s_norm.isin(UNKNOWN_LIKE)
    return df_in.copy() if include_unknown else df_in.loc[~mask_unknown_or_missing].copy()

def dq_caption(df_in: pd.DataFrame, column: str, label: str):
    if column not in df_in.columns:
        return
    s_norm = df_in[column].astype("string").str.strip().str.lower()
    mask_unknown_or_missing = df_in[column].isna() | s_norm.isin(UNKNOWN_LIKE)
    total = len(df_in)
    unknown_combined = int(mask_unknown_or_missing.sum())
    valid = int(total - unknown_combined)
    pct = (lambda n: (n / total * 100.0) if total > 0 else 0.0)
    st.caption(
        f"Data quality for **{label}** — Unknown + Missing: **{unknown_combined}** ({pct(unknown_combined):.1f}%), "
        f"Valid: **{valid}** ({pct(valid):.1f}%)."
    )

def dq_note_only(df_in: pd.DataFrame, column: str, label: str):
    """Show a compact data-quality note (Unknown + Missing merged) with NO checkbox."""
    if column not in df_in.columns:
        return
    s = df_in[column]
    s_norm = s.astype("string").str.strip().str.lower()
    mask_unknown_or_missing = s.isna() | s_norm.isin(UNKNOWN_LIKE)
    total = len(df_in)
    unknown_combined = int(mask_unknown_or_missing.sum())
    valid = int(total - unknown_combined)
    pct = lambda n: (n / total * 100.0) if total > 0 else 0.0
    st.caption(
        f"Data quality for **{label}** — Unknown + Missing: **{unknown_combined}** ({pct(unknown_combined):.1f}%), "
        f"Valid: **{valid}** ({pct(valid):.1f}%)."
    )    

def add_unknown_checkbox_and_note(
    df_in: pd.DataFrame,
    column: str,
    *,
    label: str | None = None,
    key: str | None = None,
    note_style: str = "caption",
) -> pd.DataFrame:
    label = label or column
    if column not in df_in.columns:
        st.warning(f"Column '{column}' not found; skipping Unknown filter for {label}.")
        return df_in

    s_norm = df_in[column].astype("string").str.strip().str.lower()
    mask_unknown_or_missing = df_in[column].isna() | s_norm.isin(UNKNOWN_LIKE)

    total = len(df_in)
    unknown_combined = int(mask_unknown_or_missing.sum())
    valid = int(total - unknown_combined)

    include_unknown = st.checkbox(
        f"Include 'Unknown' in {label}",
        value=False,
        key=key or f"include_unknown_{label}",
        help="When unchecked, rows with missing or 'Unknown'-like values are excluded for this section.",
    )
    filtered = df_in.copy() if include_unknown else df_in.loc[~mask_unknown_or_missing].copy()

    pct = (lambda n: (n / total * 100.0) if total > 0 else 0.0)
    note_text = (
        f"**Data quality for {label}** — Unknown + Missing: **{unknown_combined}** ({pct(unknown_combined):.1f}%), "
        f"Valid: **{valid}** ({pct(valid):.1f}%) • Analysis below **{'includes' if include_unknown else 'excludes'}** Unknown + Missing."
    )
    if note_style == "warning":
        st.warning(note_text)
    elif note_style == "caption":
        st.caption(note_text)
    else:
        st.info(note_text)

    return filtered

# Registry to capture all figures rendered in this session (for export)
# Dict: stable_name -> {"name": stable_name, "fig": go.Figure, "data": DataFrame|None, "meta": {...}}
if "export_figs" not in st.session_state or not isinstance(st.session_state["export_figs"], dict):
    st.session_state["export_figs"] = {}

def register_dataset(label: str, df: pd.DataFrame):
    if "export_figs" not in st.session_state or not isinstance(st.session_state["export_figs"], dict):
        st.session_state["export_figs"] = {}

    base = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in (label or "dataset"))
    # Prefer attaching to the last plotted figure if available
    name = st.session_state.get("last_export_name", base)

    if name in st.session_state["export_figs"]:
        st.session_state["export_figs"][name]["data"] = df.copy()
    else:
        st.session_state["export_figs"][base] = {"name": base, "fig": None, "data": df.copy(), "meta": {}}
        st.session_state["last_export_name"] = base

def _fig_to_png(fig, scale=2):
    """Return PNG bytes for a Plotly figure or None if conversion fails or bytes are not PNG."""
    try:
        # Force kaleido engine and PNG format
        b = pio.to_image(fig, format="png", scale=scale, engine="kaleido")
    except Exception as e:
        print("⚠️ to_image failed:", e)
        return None

    # Basic integrity checks: non-empty and PNG signature
    if not b or len(b) < 100:  # tiny payloads are almost certainly invalid
        print("⚠️ to_image produced empty/short bytes")
        return None
    if b[:8] != b"\x89PNG\r\n\x1a\n":
        print("⚠️ to_image did not return PNG bytes (bad magic header)")
        return None

    return b

def build_powerpoint(fig_entries: list[dict]) -> bytes:
    prs = Presentation()
    # Prefer Title Only layout (index 5) if present; fall back to Title+Content (1) or Blank (0)
    layout_idx = 5 if len(prs.slide_layouts) > 5 else (1 if len(prs.slide_layouts) > 1 else 0)
    slide_layout = prs.slide_layouts[layout_idx]

    for item in fig_entries:
        fig = item.get("fig", None)
        if fig is None:
            # Data-only entries are fine for Excel, not PPT
            continue

        # Convert to PNG safely
        img_bytes = _fig_to_png(fig, scale=2)
        if img_bytes is None:
            print("⚠️ Skipping figure (invalid PNG bytes):", item.get("name"))
            continue

        # Create slide
        slide = prs.slides.add_slide(slide_layout)

        # Title (if layout has a title placeholder)
        title_text = (item.get("meta", {}) or {}).get("title") or item.get("name") or "Chart"
        if getattr(slide.shapes, "title", None):
            slide.shapes.title.text = title_text

        # Add image (use a fresh BytesIO)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(
            img_stream,
            Inches(0.5),   # left
            Inches(1.2),   # top
            width=Inches(9.0)  # keep aspect ratio
        )

        # Footer
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9.0), Inches(0.4))
        tf = textbox.text_frame
        tf.text = f"Generated: {datetime.now():%Y-%m-%d %H:%M} | {item.get('name','')}"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    mem = io.BytesIO()
    prs.save(mem)
    mem.seek(0)
    return mem.getvalue()

def prepare_images_for_ppt(fig_entries: list[dict], scale: int = 2):
    """Return (ok, skipped) where ok contains PNG-ready items and skipped contains names that failed conversion."""
    ok, skipped = [], []
    for it in fig_entries:
        fig = it.get("fig")
        if fig is None:
            continue  # data-only items are fine for Excel, but not for PPT
        png = _fig_to_png(fig, scale=scale)
        if png is None:
            skipped.append(it.get("name", "unnamed"))
            continue
        ok.append({
            "name": it.get("name", "chart"),
            "meta": it.get("meta", {}),
            "png": png,
        })
    return ok, skipped


def build_powerpoint_from_pngs(png_entries: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    import io
    from datetime import datetime

    prs = Presentation()
    layout_idx = 5 if len(prs.slide_layouts) > 5 else (1 if len(prs.slide_layouts) > 1 else 0)
    slide_layout = prs.slide_layouts[layout_idx]

    added = 0
    for item in png_entries:
        slide = prs.slides.add_slide(slide_layout)
        title_text = (item.get("meta", {}) or {}).get("title") or item.get("name") or "Chart"
        if getattr(slide.shapes, "title", None):
            slide.shapes.title.text = title_text

        img_stream = io.BytesIO(item["png"])
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.2), width=Inches(9.0))

        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9.0), Inches(0.4))
        tf = footer.text_frame
        tf.text = f"Generated: {datetime.now():%Y-%m-%d %H:%M} | {item.get('name','')}"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        added += 1

    # If nothing got added, create a placeholder slide explaining why
    if added == 0:
        slide = prs.slides.add_slide(slide_layout)
        if getattr(slide.shapes, "title", None):
            slide.shapes.title.text = "No charts exported"
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(2.0))
        body.text_frame.text = ("No charts could be converted to images. "
                                "Check that the charts have rendered in tabs 1–8 and that Kaleido is working.")

    mem = io.BytesIO()
    prs.save(mem)
    mem.seek(0)
    return mem.getvalue()

def build_excel(fig_entries: list[dict]) -> bytes:
    mem = io.BytesIO()
    with pd.ExcelWriter(mem, engine="xlsxwriter") as writer:
        wb = writer.book
        # Style
        pct_fmt = wb.add_format({"num_format": "0.0%"})
        int_fmt = wb.add_format({"num_format": "#,##0"})

        # Write a “Contents” sheet
        contents = []
        for item in fig_entries:
            contents.append({
                "Name": item["name"],
                "Has Figure": item["fig"] is not None,
                "Has Data": item["data"] is not None
            })
        pd.DataFrame(contents).to_excel(writer, sheet_name="Contents", index=False)

        # One sheet per item (data only; images go via PPT/PDF)
        for item in fig_entries:
            if item["data"] is None:
                continue
            sheet = textwrap.shorten(item["name"], width=31, placeholder="")  # Excel sheet name limit
            df = item["data"].copy()
            df.to_excel(writer, sheet_name=sheet, index=False)
            # Optional: format numerics
            ws = writer.sheets[sheet]
            for col_idx, col in enumerate(df.columns):
                try:
                    # Apply % format if column looks like a percent series
                    if pd.api.types.is_float_dtype(df[col]) and (df[col].between(0, 1).all() or df[col].between(0, 100).all()):
                        ws.set_column(col_idx, col_idx, 14, pct_fmt)
                    else:
                        ws.set_column(col_idx, col_idx, 16, int_fmt)
                except Exception:
                    ws.set_column(col_idx, col_idx, 16)
    mem.seek(0)
    return mem.getvalue()

def build_pdf(fig_entries: list[dict]) -> bytes:
    # Simple A4 portrait PDF: one chart per page (centered)
    width, height = A4
    margin = 36  # 0.5"
    max_w = width - 2 * margin
    max_h = height - 2 * margin - 24  # leave space for title

    mem = io.BytesIO()
    c = canvas.Canvas(mem, pagesize=A4)

    for item in fig_entries:
        if item["fig"] is None:
            continue
        # Title
        title = item["meta"].get("title") or item["name"]
        c.setFont("Helvetica-Bold", 12)
        c.drawCentredString(width / 2, height - margin, title)

        # Render figure to PNG
        img_bytes = _fig_to_png(item["fig"], scale=2)
        if img_bytes:
            img = ImageReader(io.BytesIO(img_bytes))
            iw, ih = img.getSize()

            # Scale to fit
            ratio = min(max_w / iw, max_h / ih)
            w = iw * ratio
            h = ih * ratio
            x = (width - w) / 2
            y = (height - h) / 2 - 12

            c.drawImage(img, x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
            # Footer
            c.setFont("Helvetica", 8)
            c.drawRightString(width - margin, margin / 2, f"Generated {datetime.now():%Y-%m-%d %H:%M} • {item['name']}")
            c.showPage()
            c.drawImage(img, x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
        else:
            c.setFont("Helvetica", 11)
            c.drawCentredString(width/2, height/2, "Image render unavailable (install 'kaleido').")

    c.save()
    mem.seek(0)
    return mem.getvalue()

# ---- Plot helper with unique keys ----
if "plot_counter" not in st.session_state:
    st.session_state["plot_counter"] = 0

def _next_plot_key(prefix: str) -> str:
    st.session_state["plot_counter"] += 1
    return f"{prefix}_{st.session_state['plot_counter']}"

def plotly_show(fig, *, prefix: str, label: str | None = None, data_df: pd.DataFrame | None = None, **kwargs):
    # Use a unique key for Streamlit rendering only (can still use the counter)
    key = _next_plot_key(prefix)
    st.plotly_chart(fig, use_container_width=True, key=key, **kwargs)

    # Build a STABLE export name (no counter). Prefer explicit label; else prefix.
    stable_name = (label or prefix or "plot")
    safe = "".join(ch if ch.isalnum() or ch in ("-", "_") else "_" for ch in stable_name)

    # Upsert into dict so reruns overwrite instead of duplicate
    st.session_state["export_figs"][safe] = {
        "name": safe,
        "fig": fig,
        "data": (data_df.copy() if isinstance(data_df, pd.DataFrame) else st.session_state["export_figs"].get(safe, {}).get("data")),
        "meta": {"title": label or prefix}
    }

    # Remember the last export entry name (so register_dataset can attach data)
    st.session_state["last_export_name"] = safe


def safe_plot(check_df: pd.DataFrame, plot_callable):
    if isinstance(check_df, pd.DataFrame) and check_df.empty:
        st.warning("No data to display after filtering. Try including 'Unknown'.")
        return
    if isinstance(check_df, pd.Series) and check_df.empty:
        st.warning("No data to display after filtering. Try including 'Unknown'.")
        return
    plot_callable()

# ------------------------------
# Global date span (from df_full)
# ------------------------------
if "Run_Month" in df_full.columns:
    full_min = pd.to_datetime(df_full["Run_Month"], errors="coerce").min()
    full_max = pd.to_datetime(df_full["Run_Month"], errors="coerce").max()
    if "run_month_full_span" not in st.session_state:
        st.session_state["run_month_full_span"] = (full_min.date(), full_max.date())
    if "run_month_range" not in st.session_state:
        st.session_state["run_month_range"] = st.session_state["run_month_full_span"]

# ------------------------------
# Sidebar: Global buttons & date range
# ------------------------------
with st.sidebar:
    st.header("Filters")

    def select_all_filters():
        for label in UI_FILTER_LABELS:
            ms_key = _safe_key(label, "multi")
            col = _col_from_label(label)
            series = df.get(col, pd.Series([], dtype="object")).copy()
            # Avoid fillna on Categorical (can't add a new category) by casting to string first
            if pd.api.types.is_categorical_dtype(series):
                series = series.astype("string").fillna("Unknown")
            else:
                series = series.fillna("Unknown")
            st.session_state[ms_key] = sorted(series.unique().tolist())
        if "run_month_full_span" in st.session_state:
            st.session_state["run_month_range"] = st.session_state["run_month_full_span"]

    def clear_all_filters():
        for label in UI_FILTER_LABELS:
            ms_key = _safe_key(label, "multi")
            st.session_state[ms_key] = []
        if "run_month_full_span" in st.session_state:
            st.session_state["run_month_range"] = st.session_state["run_month_full_span"]

    c1, c2 = st.columns(2)
    with c1:
        st.button("✅ Select all filters", key="btn_select_all_filters", on_click=select_all_filters)
    with c2:
        st.button("🧹 Clear all filters", key="btn_clear_all_filters", on_click=clear_all_filters)

    if "Run_Month" in df.columns:
        full_min_date, full_max_date = st.session_state["run_month_full_span"]
        st.date_input(
            "Run month range",
            key="run_month_range",
            value=st.session_state["run_month_range"],
            min_value=full_min_date,
            max_value=full_max_date,
        )
        start_d, end_d = map(pd.to_datetime, st.session_state["run_month_range"])
        df = df[(df["Run_Month"] >= start_d) & (df["Run_Month"] <= end_d)]

    # Per-filter multiselects with "Select all" buttons
    sel_status   = multiselect_with_all_button("Application Status", df)
    sel_app_type = multiselect_with_all_button("Applicant Type", df)
    sel_primcat  = multiselect_with_all_button("Pri Category", df)
    sel_secncat  = multiselect_with_all_button("Sec Category", df)
    sel_country  = multiselect_with_all_button("Country", df)
    sel_senior   = multiselect_with_all_button("Seniority", df)
    sel_domain   = multiselect_with_all_button("Domain", df)
    top_k = st.number_input("Top K (for Top-X charts)", min_value=3, max_value=50, value=10, step=1)

# ------------------------------
# Apply all filters to the working df
# ------------------------------
mask = (
    apply_filter(df.get(_col_from_label("Application Status"), pd.Series(index=df.index)), sel_status) &
    apply_filter(df.get(_col_from_label("Applicant Type"),     pd.Series(index=df.index)), sel_app_type) &
    apply_filter(df.get(_col_from_label("Pri Category"),       pd.Series(index=df.index)), sel_primcat) &
    apply_filter(df.get(_col_from_label("Sec Category"),       pd.Series(index=df.index)), sel_secncat) &
    apply_filter(df.get(_col_from_label("Country"),            pd.Series(index=df.index)), sel_country) &
    apply_filter(df.get(_col_from_label("Seniority"),          pd.Series(index=df.index)), sel_senior) &
    apply_filter(df.get(_col_from_label("Domain"),             pd.Series(index=df.index)), sel_domain)
)
df_f = df[mask].copy()
st.caption(f"Filtered rows: {len(df_f):,} of {len(df):,}")

# ------------------------------
# Tabs & Visualizations
# ------------------------------
tab1, tab2, tab3, tab4, tab5, tab_6, tab_7, tab_8, tab_9 = st.tabs([
    "📈 Time Series",
    "🗺️ Geography",
    "🏷️ Programmes × Country",
    "👔 Titles & Orgs",
    "🧮 Age & Demographics",
    "🧭 Category Insights",
    "💰 Programme Cost",
    "🎯 Programme Deep Dive",
    "ℹ️ Data Preview",
])

# --- Tab 1: Time Series
with tab1:
    st.subheader("Participants over Time")
    if "Run_Month" in df_f.columns:
        ts = df_f.groupby("Run_Month").size().reset_index(name="Participants").sort_values("Run_Month")
        fig = px.line(ts, x="Run_Month", y="Participants", markers=True)
        fig.update_layout(yaxis_title="Participants", xaxis_title="Run Month")
        plotly_show(fig, prefix="tab1_participants_over_time")

    if "Programme Start Date" in df_f.columns:
        tmp = df_f.copy()
        tmp["Start_Month_Num"] = tmp["Programme Start Date"].dt.month
        tmp["Start_Quarter"]   = tmp["Programme Start Date"].dt.quarter

        col_a, col_b = st.columns(2)
        with col_a:
            mon = tmp.groupby("Start_Month_Num").size().reset_index(name="Applications")
            mon["Month"] = mon["Start_Month_Num"].map({1:"Jan",2:"Feb",3:"Mar",4:"Apr",5:"May",6:"Jun",7:"Jul",8:"Aug",9:"Sep",10:"Oct",11:"Nov",12:"Dec"})
            mon = mon.sort_values("Start_Month_Num")
            figm = px.bar(mon, x="Month", y="Applications", title="Applications by Start Month")
            plotly_show(figm, prefix="tab1_by_start_month")

        with col_b:
            q = tmp.groupby("Start_Quarter").size().reset_index(name="Applications")
            figq = px.bar(q, x="Start_Quarter", y="Applications", title="Applications by Start Quarter")
            plotly_show(figq, prefix="tab1_by_start_quarter")

# --- Tab 2: Geography 
with tab2:
    st.subheader("Geospatial: Participants by Country")

    country_col = "Country Of Residence"
    if country_col in df_f.columns:
        # DQ note only (no checkbox)
        dq_note_only(df_f, country_col, "Country")

        # Exclude Singapore toggle
        exclude_sg = st.checkbox(
            "Exclude Singapore (reduce skew)", value=False, key="geo_exclude_sg"
        )

        base = df_f.copy()
        if exclude_sg:
            base = base[base[country_col] != "Singapore"].copy()

        # --- Map: must exclude Unknown/Missing (invalid country names)
        s = base[country_col].astype("string").str.strip()
        s_norm = s.str.lower()
        mask_unknown = s.isna() | (s == "") | s_norm.isin(UNKNOWN_LIKE)
        df_map = base.loc[~mask_unknown].copy()

        if df_map.empty:
            st.info("No data to display for current filters.")
        else:
            geo = df_map.groupby(country_col).size().reset_index(name="Participants")
            size = np.sqrt(geo["Participants"].astype(float).clip(lower=0))
            size = 10 + (size / size.max()) * 30 if size.max() > 0 else np.full_like(size, 10, dtype=float)
            geo["BubbleSize"] = size

            color_values = geo["Participants"].astype(float)
            cmin = float(color_values.min())
            cmax = float(np.quantile(color_values, 0.95))
            if cmax <= cmin:
                cmax = cmin + 1.0

            fig = px.scatter_geo(
                geo,
                locations=country_col,
                locationmode="country names",
                size="BubbleSize",
                color="Participants",
                hover_name=country_col,
                hover_data={"Participants": True, "BubbleSize": False},
                color_continuous_scale="Viridis",
                projection="natural earth",
            )
            fig.update_traces(marker=dict(sizemode="area", line=dict(width=0.5, color="rgba(0,0,0,0.25)")))
            fig.update_layout(
                coloraxis_colorbar_title="Participants",
                coloraxis_cmin=cmin,
                coloraxis_cmax=cmax,
                margin=dict(l=0, r=0, t=10, b=0),
            )
            plotly_show(fig, prefix="tab2_geo_map")

        # --- Pareto bar: show Top-K countries (always excludes Unknown) ---
        st.markdown("**Pareto of Countries (Top K)**")
        if df_f.empty:
            st.info("No countries to display for the current filters.")
        else:
            pareto_base = df_f.copy()

            # Exclude only Singapore, case-insensitive
            if exclude_sg:
                pareto_base = pareto_base[~_norm_str(pareto_base[country_col]).eq("singapore")].copy()

            # Always drop Unknown / Missing
            s_norm = _norm_str(pareto_base[country_col])
            mask_unknown = pareto_base[country_col].isna() | s_norm.isin(UNKNOWN_LIKE)
            s = pareto_base.loc[~mask_unknown, country_col]

            if s.empty:
                st.info("No valid countries to display for the current filters.")
                st.stop()

            # Counts for the FULL shown universe (after SG filter, Unknown removed)
            counts_df = s.value_counts(dropna=False).reset_index()
            counts_df.columns = [country_col, "Participants"]

            # This is the denominator for percentages and for the "Top K countries account for ..." note
            total_universe = int(counts_df["Participants"].sum())

            # Top-K slice (non-Unknown by construction)
            final_df = counts_df.nlargest(int(top_k), "Participants").copy()

            # Percent labels and caption share are relative to the FULL shown universe
            final_df["Share_%"] = (final_df["Participants"] / total_universe * 100.0) if total_universe > 0 else 0.0

            fig_bar = px.bar(
                final_df,
                x=country_col,
                y="Participants",
                title=f"Top {int(top_k)} Countries by Participants",
                text=final_df["Share_%"].round(1).astype(str) + "%",
            )
            fig_bar.update_traces(textposition="outside", cliponaxis=False)
            fig_bar.update_layout(xaxis_tickangle=-45, yaxis_title="Participants", xaxis_title="Country")

            # ---- attach data for Excel export
            export_df_tab2 = final_df[[country_col, "Participants", "Share_%"]].copy()
            export_df_tab2.rename(columns={country_col: "Country"}, inplace=True)

            plotly_show(
                fig_bar,
                prefix="tab2_geo_pareto",
                label="Tab2 — Top Countries by Participants (Pareto)",
                data_df=export_df_tab2
            )
            sg_note = " (Singapore excluded)" if exclude_sg else ""
            st.caption(
                f"Total participants shown: {total_universe:,}{sg_note}. "
                f"Top {int(top_k)} countries account for {final_df['Share_%'].sum():.1f}% of the shown total."
            )



# --- Tab 3: Programmes × Country (both heatmaps share one %/raw toggle) ---
with tab3:
    st.subheader("Top Programmes & Country Breakdown")

    prog_col    = "Truncated Programme Name"
    country_col = "Country Of Residence"

    if (prog_col in df_f.columns) and (country_col in df_f.columns):
        # DQ note only (no checkbox)
        dq_note_only(df_f, country_col, "Country (Heatmaps)")

        # Explain ordering logic (applies to both heatmaps)
        st.caption("Rows in both heatmaps are ordered by total participants after filtering (Unknown removed; Singapore excluded if selected).")

        # Shared controls for both heatmaps
        exclude_sg_tab3 = st.checkbox(
            "Exclude Singapore (reduce skew)", value=False, key="pc_exclude_sg"
        )
        show_pct_tab3 = st.toggle(
            "Show % (vs raw counts)", value=True, key="tab3_show_pct"
        )

        # Base (apply SG exclusion consistently for both HMs)
        base = df_f.copy()
        if exclude_sg_tab3:
            norm_cty = base[country_col].astype("string").str.strip().str.casefold()
            base = base.loc[~norm_cty.eq("singapore")].copy()

        # -------------------------------
        # Heatmap 1: Programme × Country
        # -------------------------------
        st.markdown("### Heatmap 1 — Programme × Country")

        # Remove Unknown/Missing countries for interpretability
        s = base[country_col].astype("string").str.strip()
        s_norm = s.str.lower()
        mask_unknown = s.isna() | (s == "") | s_norm.isin(UNKNOWN_LIKE)
        base_hm1 = base.loc[~mask_unknown].copy()

        if base_hm1.empty:
            st.info("No valid Country data for Heatmap 1 after filtering.")
        else:
            # Limit rows/cols to Top-K by volume after filtering
            top_progs = (
                base_hm1[prog_col].value_counts()
                .nlargest(int(top_k))
                .index
                .tolist()
            )
            agg = (
                base_hm1[base_hm1[prog_col].isin(top_progs)]
                .groupby([prog_col, country_col])
                .size()
                .reset_index(name="Participants")
            )
            top_countries = (
                agg.groupby(country_col)["Participants"].sum()
                .nlargest(int(top_k))
                .index
                .tolist()
            )
            agg = agg[agg[country_col].isin(top_countries)]

            # Pivot to matrix
            hm1_counts = (
                agg.pivot(index=prog_col, columns=country_col, values="Participants")
                .fillna(0)
            )

            # Order rows by total participants (desc) AFTER filtering
            hm1_counts = hm1_counts.loc[
                hm1_counts.sum(axis=1).sort_values(ascending=False).index
            ]

            # Choose matrix + labels + hover per mode
            if show_pct_tab3:
                total_sum = hm1_counts.values.sum()
                hm1_pct = (hm1_counts / total_sum * 100).round(2) if total_sum > 0 else hm1_counts * 0
                Z = hm1_pct.values
                x_labels, y_labels = hm1_pct.columns, hm1_pct.index
                color_label = "Share of total (%)"
                title_suffix = "(% of total)"
                text_fmt = ".2f"
                hover_tmpl = (
                    "Programme: %{y}<br>"
                    "Country: %{x}<br>"
                    "Share of total: %{z:.2f}%"
                    "<extra></extra>"
                )
            else:
                Z = hm1_counts.values
                x_labels, y_labels = hm1_counts.columns, hm1_counts.index
                color_label = "Participants"
                title_suffix = "(raw)"
                text_fmt = "d"
                hover_tmpl = (
                    "Programme: %{y}<br>"
                    "Country: %{x}<br>"
                    "Participants: %{z:.0f}"
                    "<extra></extra>"
                )

            fig_hm1 = px.imshow(
                Z,
                x=x_labels,
                y=y_labels,
                color_continuous_scale="Viridis",
                aspect="auto",
                labels=dict(x="Country Of Residence", y="Programme", color=color_label),
                text_auto=text_fmt,
                title=f"Programme × Country {title_suffix}",
            )
            fig_hm1.update_traces(hovertemplate=hover_tmpl)
            fig_hm1.update_layout(xaxis_title="Country Of Residence", yaxis_title="Programme (Anon)")
            
            # ---- attach data for Excel export (both raw counts and overall %)
            df_export_hm1 = agg.copy()
            total_hm1 = df_export_hm1["Participants"].sum()  # ✅ Safely recompute total here

            df_export_hm1["Overall_%"] = (
                df_export_hm1["Participants"] / total_hm1 * 100.0
            ).round(2) if total_hm1 > 0 else 0.0

            df_export_hm1.rename(columns={
                prog_col: "Programme",
                country_col: "Country"
            }, inplace=True)

            plotly_show(
                fig_hm1,
                prefix="tab3_prog_country_heatmap",
                label="Tab3 — Programme × Country (Overall %)",
                data_df=df_export_hm1
            )

        # -------------------------------------------------------------
        # Heatmap 2: Top Countries × Primary Category (row % or raw)
        # -------------------------------------------------------------
        st.markdown(f"### Heatmap 2 — Top {top_k} Countries × Primary Category")

        cat_col = "Primary Category"
        if cat_col not in df_f.columns:
            st.info("‘Primary Category’ column not found for Heatmap 2.")
        else:
            # Remove Unknown/Missing countries; SG already excluded via base
            s2 = base[country_col].astype("string").str.strip()
            s2_norm = s2.str.lower()
            mask_unknown2 = s2.isna() | (s2 == "") | s2_norm.isin(UNKNOWN_LIKE)
            base_hm2 = base.loc[~mask_unknown2].copy()

            if base_hm2.empty:
                st.info("No valid Country data for Heatmap 2 after filtering.")
            else:
                # Pick Top-K countries by total participants AFTER filtering
                top_countries_hm2 = (
                    base_hm2[country_col].value_counts()
                    .nlargest(int(top_k))
                    .index
                    .tolist()
                )
                df_top_cty = base_hm2[base_hm2[country_col].isin(top_countries_hm2)].copy()

                # Build counts matrix: rows=Country, cols=Primary Category
                agg_cat = (
                    df_top_cty
                    .groupby([country_col, cat_col])
                    .size()
                    .reset_index(name="Participants")
                )
                hm2_counts = (
                    agg_cat
                    .pivot(index=country_col, columns=cat_col, values="Participants")
                    .fillna(0)
                )

                # Order rows by total participants (desc) AFTER filtering
                hm2_counts = hm2_counts.loc[
                    hm2_counts.sum(axis=1).sort_values(ascending=False).index
                ]

                if show_pct_tab3:
                    # Row % version (each row sums to 100)
                    row_sums = hm2_counts.sum(axis=1).replace(0, np.nan)
                    hm2_pct = (hm2_counts.div(row_sums, axis=0) * 100).round(2).fillna(0)
                    Z2 = hm2_pct.values
                    x2, y2 = hm2_pct.columns, hm2_pct.index
                    color_label2 = "Row %"
                    title_suffix2 = "(row %)"
                    text_fmt2 = ".2f"
                    hover_tmpl2 = (
                        "Country: %{y}<br>"
                        "Primary Category: %{x}<br>"
                        "Row share: %{z:.2f}%"
                        "<extra></extra>"
                    )
                else:
                    Z2 = hm2_counts.values
                    x2, y2 = hm2_counts.columns, hm2_counts.index
                    color_label2 = "Participants"
                    title_suffix2 = "(raw)"
                    text_fmt2 = "d"
                    hover_tmpl2 = (
                        "Country: %{y}<br>"
                        "Primary Category: %{x}<br>"
                        "Participants: %{z:.0f}"
                        "<extra></extra>"
                    )

                fig_cat = px.imshow(
                    Z2,
                    x=x2,
                    y=y2,
                    color_continuous_scale="Viridis",
                    aspect="auto",
                    labels=dict(x="Primary Category", y="Country Of Residence", color=color_label2),
                    text_auto=text_fmt2,
                    title=f"For each country: distribution {title_suffix2}",
                )
                fig_cat.update_traces(hovertemplate=hover_tmpl2)
                fig_cat.update_layout(xaxis_title="Primary Category", yaxis_title="Country Of Residence")
                
                # ---- attach data for Excel export (raw counts + row %)
                df_export_hm2 = agg_cat.copy()
                # compute row totals and row percentages
                row_totals = df_export_hm2.groupby(country_col)["Participants"].transform("sum")
                df_export_hm2["Row_%"] = (df_export_hm2["Participants"] / row_totals * 100.0).round(2)
                df_export_hm2.rename(columns={
                    country_col: "Country",
                    "Primary Category": "Primary Category",
                    "Participants": "Participants"
                }, inplace=True)

                plotly_show(
                    fig_cat,
                    prefix="tab3_country_primarycat_heatmap",
                    label="Tab3 — Country × Primary Category (Row %)",
                    data_df=df_export_hm2
                )
    else:
        st.info("Required columns not found: ensure ‘Truncated Programme Name’ and ‘Country Of Residence’ exist in the dataset.")

# --- Tab 4: Titles & Organisations
with tab4:
    st.subheader("Top Job Titles & Organisations")

    if "Job Title Clean" in df_f.columns:
        df_title = add_unknown_checkbox_and_note(
            df_f, "Job Title Clean", label="Job Title", key="job_title_tab4", note_style="caption"
        )

        # When Unknowns are included, coalesce blank/variants into literal "Unknown"
        s_titles = df_title["Job Title Clean"]
        if "Unknown" in df_title["Job Title Clean"].astype("string").str.lower().unique() or df_title["Job Title Clean"].isna().any():
            s_titles = coalesce_unknown(s_titles)

        top_titles = (
            s_titles.value_counts(dropna=False)
            .head(top_k)
            .reset_index()
        )
        top_titles.columns = ["Job Title", "Participants"]

        safe_plot(
            top_titles,
            lambda: plotly_show(
                px.bar(
                    top_titles, x="Participants", y="Job Title",
                    orientation="h", title=f"Top {top_k} Job Titles"
                ),
                prefix="tab4_top_titles"
            )
        )

    org_col = "Organisation Name: Organisation Name"
    if org_col in df_f.columns:
        df_org = add_unknown_checkbox_and_note(df_f, org_col, label="Organisation", key="orgs", note_style="caption")
        top_orgs = df_org[org_col].value_counts().nlargest(top_k).reset_index()
        top_orgs.columns = ["Organisation", "Participants"]
        safe_plot(top_orgs, lambda: plotly_show(px.bar(top_orgs, x="Participants", y="Organisation", orientation="h", title=f"Top {top_k} Organisations"), prefix="tab4_top_orgs", theme="streamlit"))

    if "Domain" in df_f.columns:
        df_dom = add_unknown_checkbox_and_note(
            df_f,
            "Domain",
            label="Domain",
            key="domain_tab4",         
            note_style="caption"
        )

        # Read the checkbox state that add_unknown_checkbox_and_note created
        include_unknown_domain = bool(st.session_state.get("domain_tab4", False))

        # If the checkbox is UNCHECKED, also hide literal "Others"
        # 'Others' which is a valid classified category based on BERTopic results will be included in the 'Valid' count.
        if not include_unknown_domain:
            mask_others = (
                df_dom["Domain"].astype("string").str.strip().str.lower() == "others"
            )
            df_dom = df_dom.loc[~mask_others].copy()

        # Now compute the chart from the (optionally) trimmed df_dom
        top_domains = (
            df_dom["Domain"]
            .value_counts(dropna=False)
            .nlargest(int(top_k))
            .reset_index()
        )
        top_domains.columns = ["Domain", "Participants"]

        # Additional data-quality clarity note for Domain
        st.caption(
            "Note: **'Others'** is a valid derived category from clustering — "
            "it is hidden by default for clearer insights, but **still counted as valid data** "
            "in the data quality stats above."
        )

    safe_plot(
        top_domains,
        lambda: plotly_show(
            px.bar(
                top_domains,
                x="Participants",
                y="Domain",
                orientation="h",
                title=f"Top {int(top_k)} Domains",
            ),
            prefix="tab4_top_domains",
            theme="streamlit",
        ),
    )

    if "Seniority" in df_f.columns:
        df_sen = add_unknown_checkbox_and_note(df_f, "Seniority", key="seniority", note_style="caption")
        sen = df_sen["Seniority"].value_counts().reset_index()
        sen.columns = ["Seniority", "Participants"]
        safe_plot(sen, lambda: plotly_show(px.bar(sen, x="Seniority", y="Participants", title="Participants by Seniority"), prefix="tab4_seniority"))
    
    

# --- Tab 5: Age & Demographics
with tab5:
    st.subheader("Demographics")

    if "Age_Group" in df_f.columns:
        # This renders the checkbox + DQ note and returns the filtered df
        df_age = add_unknown_checkbox_and_note(
            df_f, "Age_Group", label="Age Group", key="agegroup_tab5", note_style="caption"
        )

        # Read the checkbox state without modifying the helper
        include_unknown_age = bool(st.session_state.get("agegroup_tab5", False))

        s_raw = df_age["Age_Group"].astype("string")

        if include_unknown_age:
            # Merge any blanks/NA/unknown-like into literal "Unknown"
            s = coalesce_unknown(s_raw)
            order = ["<35", "35–44", "45–54", "55–64", "65+", "Unknown"]
        else:
            # Ensure Unknown never appears when checkbox is OFF
            s = s_raw[s_raw.ne("Unknown")]
            order = ["<35", "35–44", "45–54", "55–64", "65+"]

        agec = (
            s.value_counts(dropna=False)
             .reindex(order)
             .fillna(0)
             .rename_axis("Age Group")
             .reset_index(name="Participants")
        )

        safe_plot(
            agec,
            lambda: plotly_show(
                px.bar(agec, x="Age Group", y="Participants", title="Participants by Age Group"),
                prefix="tab5_agegroup_bar"
            )
        )

    if "Gender" in df_f.columns:
        df_gender = add_unknown_checkbox_and_note(df_f, "Gender", key="gender_tab5", note_style="caption")
        gender = df_gender["Gender"].value_counts().reset_index()
        gender.columns = ["Gender", "Participants"]
        safe_plot(gender, lambda: plotly_show(px.pie(gender, names="Gender", values="Participants", title="Gender Split"), prefix="tab5_gender_pie"))

# --- Tab 6: Category Insights
with tab_6:
    st.subheader("Category Insights")
    sub_age, sub_country = st.tabs(["📊 Age Distribution per Category", "🌍 Country Distribution per Category"])

    with sub_age:
        st.markdown("##### Age Distribution per Category")
        cat_type = st.radio(
            "Choose category type:",
            ["Primary Category", "Secondary Category"],
            key="age_cat_type",
            horizontal=True
        )
        cat_col = cat_type

        if (cat_col in df_f.columns) and ("Age_Group" in df_f.columns):
            cat_values = (
                df_f[cat_col].astype("string").fillna("Unknown").replace({"": "Unknown"}).unique().tolist()
            )
            # Put Unknown last
            cat_values = [v for v in sorted(cat_values) if v != "Unknown"] + (
                ["Unknown"] if "Unknown" in cat_values else []
            )
            selected_cat = st.selectbox(f"Select {cat_type}:", cat_values, key="age_cat_select")

            subset = df_f[df_f[cat_col].astype("string").fillna("Unknown") == selected_cat].copy()
            if subset.empty:
                st.info("No rows for this selection.")
            else:
                # Checkbox + DQ note for Age Group in this selection
                sub_age_df = add_unknown_checkbox_and_note(
                    subset, "Age_Group", label="Age Group (this selection)",
                    key="age_dist_sub", note_style="caption"
                )

                # Coalesce Unknown-like + Missing -> "Unknown" BEFORE counting
                s = coalesce_unknown(sub_age_df["Age_Group"])

                # Build % distribution
                dist = (s.value_counts(normalize=True, dropna=False) * 100.0).reset_index()
                dist.columns = ["Age Group", "Percentage"]

                # Order buckets; keep Unknown last when present
                order_full = ["<35", "35–44", "45–54", "55–64", "65+", "Unknown"]
                # Only keep groups that exist after filtering
                present = [g for g in order_full if g in dist["Age Group"].values]
                dist = dist.set_index("Age Group").reindex(present).reset_index()

                # Guard against empty after filtering
                if dist.empty:
                    st.info("No Age Group data to display for this selection.")
                else:
                    text_labels = dist["Percentage"].round(1).astype(str) + "%"
                    fig = px.bar(
                        dist,
                        x="Age Group",
                        y="Percentage",
                        title=f"Age Distribution (%) – {cat_type}: {selected_cat}",
                        text=text_labels,
                    )
                    fig.update_traces(textposition="outside", cliponaxis=False)
                    ymax = min(100.0, float(dist["Percentage"].max()) + 10.0)
                    fig.update_layout(yaxis_range=[0, ymax])
                    plotly_show(fig, prefix="tab6_age_dist_by_cat")
        else:
            st.info("Required columns not found: ensure ‘Age_Group’ and the selected category column exist.")

    with sub_country:
        st.markdown("##### Country Distribution per Category")

        cat_type = st.radio(
            "Choose category type:",
            ["Primary Category", "Secondary Category"],
            key="country_cat_type",
            horizontal=True
        )
        cat_col = cat_type
        country_col = "Country Of Residence"

        # Exclude SG toggle (same behavior as Tab 2 Pareto)
        exclude_sg_tab6 = st.checkbox(
            "Exclude Singapore (reduce skew)", value=False, key="tab6_exclude_sg"
        )

        if (cat_col in df_f.columns) and (country_col in df_f.columns):
            # Category selector (keep Unknown at the end, just like elsewhere)
            cat_values = (
                df_f[cat_col].astype("string").fillna("Unknown").replace({"": "Unknown"}).unique().tolist()
            )
            cat_values = [v for v in sorted(cat_values) if v != "Unknown"] + (["Unknown"] if "Unknown" in cat_values else [])
            selected_cat = st.selectbox(f"Select {cat_type}:", cat_values, key="country_cat_select")

            subset = df_f[df_f[cat_col].astype("string").fillna("Unknown") == selected_cat].copy()
            if subset.empty:
                st.info("No rows for this selection.")
            else:
                # Same DQ note pattern as Tab 2
                dq_note_only(subset, country_col, "Country (this selection)")

                # Apply Exclude-SG first, like Tab 2 (case-insensitive)
                if exclude_sg_tab6:
                    subset = subset[~_norm_str(subset[country_col]).eq("singapore")].copy()

                # Always drop Unknown/Missing, like Tab 2 Pareto
                s_norm = _norm_str(subset[country_col])
                mask_unknown = subset[country_col].isna() | s_norm.isin(UNKNOWN_LIKE)
                sub_valid = subset.loc[~mask_unknown].copy()

                if sub_valid.empty:
                    st.info("No valid countries to display after removing Unknown (and Singapore, if excluded).")
                else:
                    # Counts over the FULL shown universe (after SG filter, Unknown removed)
                    counts_df = sub_valid[country_col].value_counts(dropna=False).reset_index()
                    counts_df.columns = [country_col, "Participants"]

                    # Denominator for percentages and caption share — same as Tab 2
                    total_universe = int(counts_df["Participants"].sum())

                    # Top-K slice
                    final_df = counts_df.nlargest(int(top_k), "Participants").copy()
                    final_df["Share_%"] = (final_df["Participants"] / total_universe * 100.0) if total_universe > 0 else 0.0

                    # Plot with % labels (same style as Tab 2 Pareto)
                    fig = px.bar(
                        final_df,
                        x=country_col,
                        y="Participants",
                        title=f"Top {int(top_k)} Countries by Participants — {cat_type}: {selected_cat}",
                        text=final_df["Share_%"].round(1).astype(str) + "%",
                    )
                    fig.update_traces(textposition="outside", cliponaxis=False)
                    fig.update_layout(xaxis_tickangle=-45, yaxis_title="Participants", xaxis_title="Country")
                    plotly_show(fig, prefix="tab6_country_dist_by_cat_like_tab2")

                    sg_note = " (Singapore excluded)" if exclude_sg_tab6 else ""
                    st.caption(
                        f"Total participants shown: {total_universe:,}{sg_note}. "
                        f"Top {int(top_k)} countries account for {final_df['Share_%'].sum():.1f}% of the shown total."
                    )
        else:
            st.info("Required columns not found: make sure ‘Country Of Residence’ and category columns exist in the dataset.")

# --- Tab 7: Programme Cost
with tab_7:
    st.subheader("Programme Cost Analysis")

    required_cols = ["Programme Cost", "Truncated Programme Name", "Run_Month"]
    if not all(col in df_f.columns for col in required_cols):
        st.warning("The required columns ('Programme Cost', 'Truncated Programme Name', 'Run_Month') are not available in the data.")
    else:
        df_cost = df_f.copy()
        df_cost['Programme Cost'] = pd.to_numeric(df_cost['Programme Cost'], errors='coerce')
        df_cost.dropna(subset=['Programme Cost'], inplace=True)

        if df_cost.empty:
            st.info("No data with valid programme costs found for the current filters.")
        else:
            st.markdown("##### Enrolment Volume vs. Programme Cost")
            grouped = df_cost.groupby('Truncated Programme Name').agg(
                enrolment_volume=('Truncated Programme Name', 'size'),
                programme_cost=('Programme Cost', 'first')
            ).reset_index()

            fig_scatter = px.scatter(
                grouped,
                x='programme_cost',
                y='enrolment_volume',
                title='Enrolment Volume vs. Programme Cost',
                labels={'programme_cost': 'Programme Cost ($)', 'enrolment_volume': 'Total Enrolments'},
                hover_data=['Truncated Programme Name']
            )
            fig_scatter.update_traces(marker=dict(size=12, opacity=0.7, line=dict(width=1, color='DarkSlateGrey')))
            plotly_show(fig_scatter, prefix="tab7_cost_vs_enrolment")

            st.divider()
            st.markdown("##### Monthly Revenue Trend")
            monthly_revenue = df_cost.groupby(df_cost['Run_Month'].dt.to_period('M'))['Programme Cost'].sum().reset_index()
            monthly_revenue.rename(columns={'Programme Cost': 'Total_Revenue'}, inplace=True)
            monthly_revenue['Run_Month'] = monthly_revenue['Run_Month'].dt.to_timestamp()
            monthly_revenue = monthly_revenue.sort_values("Run_Month")

            fig_trend = px.line(monthly_revenue, x='Run_Month', y='Total_Revenue', title='Monthly Revenue Trend', labels={'Run_Month': 'Month', 'Total_Revenue': 'Total Revenue ($)'}, markers=True)
            fig_trend.update_layout(yaxis_title="Total Revenue ($)", xaxis_title="Month")
            plotly_show(fig_trend, prefix="tab7_monthly_revenue")

            st.divider()
            st.markdown("##### Top K Countries by Total Revenue")

            if "Country Of Residence" in df_cost.columns:
                top_countries = (
                    df_cost.groupby("Country Of Residence")["Programme Cost"]
                    .sum()
                    .nlargest(top_k)
                    .reset_index()
                )
                top_countries.columns = ["Country Of Residence", "Total Revenue"]

                fig2 = px.bar(
                    top_countries,
                    x="Total Revenue",
                    y="Country Of Residence",
                    orientation="h",
                    title=f"Top {top_k} Countries by Total Revenue" 
                )
                st.plotly_chart(fig2, use_container_width=True, theme="streamlit") 

# --- Tab 8: Programme Deep Dive
with tab_8:
    st.subheader("Programme Deep Dive")
    include_unknown_deep = st.checkbox("Include 'Unknown' in Deep Dive visuals", value=False, key="dd_include_unknown")

    prog_col = "Truncated Programme Name"
    if prog_col not in df_f.columns:
        st.info("Programme column not found in the filtered data.")
        st.stop()

    progs = (df_f[prog_col].dropna().astype(str).sort_values().unique().tolist()) if not df_f.empty else []
    if not progs:
        st.info("No programmes available under current filters.")
        st.stop()

    sel_prog = st.selectbox("Select a programme", progs, index=0, key="prog_dd_select")
    p = df_f[df_f[prog_col] == sel_prog].copy()
    if p.empty:
        st.info("No rows for this programme with current filters.")
        st.stop()

    c1, c2, c3, c4, c5 = st.columns(5)
    with c1:
        st.metric("Participants", f"{len(p):,}")
    with c2:
        st.metric("Unique Runs", int(p.get("Truncated Programme Run", pd.Series()).nunique() if "Truncated Programme Run" in p.columns else 0))
    with c3:
        st.metric("Countries", int(p.get("Country Of Residence", pd.Series()).nunique() if "Country Of Residence" in p.columns else 0))
    with c4:
        st.metric("Median Age", f"{p['Age'].median():.0f}" if "Age" in p.columns and p["Age"].notna().any() else "—")
    with c5:
        st.metric("Cost", p["Programme Cost"].unique() if p["Programme Cost"].notna().any() else "Unknown")
    date_min = pd.to_datetime(p.get("Programme Start Date")).min() if "Programme Start Date" in p.columns else pd.NaT
    date_max = pd.to_datetime(p.get("Programme End Date")).max() if "Programme End Date" in p.columns else pd.NaT
    if pd.notna(date_min) and pd.notna(date_max):
        st.caption(f"**Programme Date Range:** {date_min:%A, %d %B %Y} → {date_max:%A, %d %B %Y}")
        st.divider()

    if "Run_Month" in p.columns:
        ts = p.groupby("Run_Month").size().reset_index(name="Participants").sort_values("Run_Month")
        fig_ts = px.line(ts, x="Run_Month", y="Participants", markers=True, title="Participants over Time (by Run Month)")
        fig_ts.update_layout(yaxis_title="Participants", xaxis_title="Run Month")
        fig_ts.update_xaxes(tickformat="%b %Y")
        fig_ts.update_traces(hovertemplate="Run Month=%{x|%b %Y}<br>Participants=%{y}<extra></extra>")
        plotly_show(fig_ts, prefix="tab8_prog_ts")

    colL, colR = st.columns(2)
    with colL:
        if "Application Status" in p.columns:
            p_status = filter_unknown_no_ui(p, "Application Status", include_unknown_deep)
            dq_caption(p, "Application Status", "Application Status")
            status = p_status["Application Status"].value_counts().reset_index()
            status.columns = ["Application Status", "Count"]
            safe_plot(status, lambda: plotly_show(px.bar(status, x="Application Status", y="Count", title="Application Status Breakdown", text="Count").update_traces(textposition="outside", cliponaxis=False).update_layout(xaxis_tickangle=-30), prefix="tab8_status_bar"))

        if "Gender" in p.columns:
            p_gender = filter_unknown_no_ui(p, "Gender", include_unknown_deep)
            dq_caption(p, "Gender", "Gender")
            gender = p_gender["Gender"].value_counts().reset_index()
            gender.columns = ["Gender", "Participants"]
            safe_plot(gender, lambda: plotly_show(px.pie(gender, names="Gender", values="Participants", title="Gender Split"), prefix="tab8_gender_pie"))

    with colR:
        if "Country Of Residence" in p.columns:
            p_cty = filter_unknown_no_ui(p, "Country Of Residence", include_unknown_deep)
            dq_caption(p, "Country Of Residence", "Country")
            s = p_cty["Country Of Residence"]
            if include_unknown_deep:
                s = coalesce_unknown(s)

            if s.empty:
                st.info("No country data to display for this selection.")
            else:
                ctry_counts = s.value_counts(dropna=False).reset_index()
                ctry_counts.columns = ["Country", "Participants"]
                unknown_row = ctry_counts[ctry_counts["Country"] == "Unknown"]
                non_unknown = ctry_counts[ctry_counts["Country"] != "Unknown"]
                top_non_unknown = non_unknown.nlargest(top_k, "Participants")
                ctry_top = pd.concat([top_non_unknown, unknown_row], ignore_index=True).drop_duplicates(subset=["Country"])

                order = ctry_top.sort_values("Participants", ascending=False)["Country"].tolist()
                if "Unknown" in order:
                    order = [c for c in order if c != "Unknown"] + ["Unknown"]
                    ctry_top["Country"] = pd.Categorical(ctry_top["Country"], categories=order, ordered=True)
                    ctry_top = ctry_top.sort_values("Country")

                safe_plot(ctry_top, lambda: plotly_show(px.bar(ctry_top, x="Participants", y="Country", orientation="h", title=f"Top {top_k} Countries (Participants)"), prefix="tab8_top_countries"))

        org_col = "Organisation Name: Organisation Name"
        if org_col in p.columns:
            p_org = filter_unknown_no_ui(p, org_col, include_unknown_deep)
            dq_caption(p, org_col, "Organisation")
            orgs = p_org[org_col].value_counts().reset_index()
            orgs.columns = ["Organisation", "Participants"]
            orgs_top = orgs.head(top_k)
            safe_plot(orgs_top, lambda: plotly_show(px.bar(orgs_top, x="Participants", y="Organisation", orientation="h", title=f"Top {top_k} Organisations (Participants)"), prefix="tab8_top_orgs"))

    colA, colB = st.columns(2)
    with colA:
        if "Seniority" in p.columns:
            p_sen = filter_unknown_no_ui(p, "Seniority", include_unknown_deep)
            dq_caption(p, "Seniority", "Seniority")
            sen = p_sen["Seniority"].value_counts().reset_index()
            sen.columns = ["Seniority", "Participants"]
            safe_plot(sen, lambda: plotly_show(px.bar(sen, x="Seniority", y="Participants", title="Seniority Mix"), prefix="tab8_seniority"))

    with colB:
        if "Age_Group" in p.columns:
            # Apply Unknown filter based on deep dive toggle
            p_age = filter_unknown_no_ui(p, "Age_Group", include_unknown_deep)
            dq_caption(p, "Age_Group", "Age Group")

            s = p_age["Age_Group"].astype("string")

            # If include_unknown is ON, merge Unknown-like → "Unknown"
            if include_unknown_deep:
                s = coalesce_unknown(s)

            # Build counts after filtering
            age_counts = s.value_counts(dropna=False)

            # Define ordering and respect filtered set
            order_full = ["<35","35–44","45–54","55–64","65+","Unknown"]
            present = [g for g in order_full if g in age_counts.index]

            if not present:
                st.info("No Age Group data to display for this programme.")
            else:
                age_df = (
                    age_counts.reindex(present)
                    .fillna(0)
                    .rename_axis("Age Group")
                    .reset_index(name="Participants")
                )

                fig = px.bar(
                    age_df,
                    x="Age Group",
                    y="Participants",
                    title="Age Group Distribution",
                    text="Participants"
                )
                fig.update_traces(textposition="outside", cliponaxis=False)
                plotly_show(fig, prefix="tab8_agegroup")


    st.markdown("##### Runs for this Programme")
    run_cols = [c for c in ["Truncated Programme Run", "Programme Start Date", "Programme End Date", "Country Of Residence", "Application Status"] if c in p.columns]
    st.dataframe(p.sort_values(["Run_Month","Programme Start Date"]).loc[:, run_cols].head(500), use_container_width=True, hide_index=True)

    st.download_button(
        f"Download '{sel_prog}' rows (CSV)",
        data=p.to_csv(index=False).encode("utf-8-sig"),
        file_name=f"{sel_prog[:40].replace(' ','_')}_export.csv",
        mime="text/csv"
    )

# --- Tab 9: Data Preview
with tab_9:
    st.subheader("Filtered Data Preview")
    preview_cols = [c for c in [
        "Application ID", "Contact ID", "Application Status", "Applicant Type",
        "Organisation Name: Organisation Name", "Job Title Clean", "Seniority",
        "Truncated Programme Name", "Truncated Programme Run", "Primary Category", "Secondary Category",
        "Programme Start Date", "Programme End Date", "Run_Month", "Run_Month_Label",
        "Gender", "Age", "Country Of Residence", "Domain", "Programme Cost"
    ] if c in df_f.columns]

    st.dataframe(df_f.sort_values("Run_Month").loc[:, preview_cols].head(500), use_container_width=True, hide_index=True)
    st.download_button("Download filtered CSV", data=df_f.to_csv(index=False).encode("utf-8-sig"), file_name="filtered_export.csv", mime="text/csv")

tab_export, = st.tabs(["📤 Exports"])
with tab_export:
    st.subheader("Export Dashboard Insights")
    st.caption("Exports reflect the **current filters**. Visit tabs 1–8 first so charts render and get captured here.")

    registry = st.session_state.get("export_figs", {})
    items = [v for v in registry.values() if (v.get("fig") is not None or v.get("data") is not None)]

    st.info(
        f"Charts captured: **{sum(1 for v in registry.values() if v.get('fig') is not None)}** • "
        f"Tables captured: **{sum(1 for v in registry.values() if v.get('data') is not None)}**"
    )

    export_choice = st.selectbox(
        "Choose export format",
        ["PowerPoint (.pptx)", "Excel (.xlsx)", "PDF (.pdf)", "All (ZIP with PPTX+XLSX+PDF)"],
        index=0,
        key="export_choice_tab10",
    )

    # Build files
    if st.button("Build export files"):
        items = [x for x in st.session_state["export_figs"] if (x.get("fig") is not None or x.get("data") is not None)]
        if not items:
            st.warning("No charts/tables captured yet. Navigate tabs 1–8 so visuals render first.")
        else:
            built = []

            # Pre-render PNGs once so we know exactly how many valid charts we have
            png_ok, png_skipped = prepare_images_for_ppt(items, scale=2)

            if png_skipped:
                st.warning("Skipped charts (failed image conversion): " + ", ".join(png_skipped[:10]) + ("..." if len(png_skipped) > 10 else ""))

            if want_ppt:
                ppt_bytes = build_powerpoint_from_pngs(png_ok)
                built.append(("EE_Insights.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", ppt_bytes))

            if want_xls:
                xls_bytes = build_excel(items)  # Excel uses dataframes; keep as-is
                built.append(("EE_Insights.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", xls_bytes))

            if want_pdf:
                # PDF still renders from figures directly; okay to keep your existing build_pdf(items)
                pdf_bytes = build_pdf(items)
                built.append(("EE_Insights.pdf", "application/pdf", pdf_bytes))

            if not built:
                st.info("No formats selected.")
            elif add_zip and len(built) > 1:
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
                    for fname, _mime, b in built:
                        zf.writestr(fname, b)
                buf.seek(0)
                st.download_button("⬇️ Download exports.zip", data=buf.getvalue(), file_name="EE_exports.zip", mime="application/zip")
            else:
                for fname, mime, b in built:
                    st.download_button(f"⬇️ Download {fname}", data=b, file_name=fname, mime=mime)

